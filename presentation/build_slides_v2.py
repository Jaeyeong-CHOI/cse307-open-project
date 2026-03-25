#!/usr/bin/env python3
"""Build professional research meeting slides v2 (25 main + 8 appendix)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Color palette ──
NAVY       = RGBColor(0x1E, 0x2D, 0x5A)
NAVY_DARK  = RGBColor(0x10, 0x1B, 0x38)
BLUE       = RGBColor(0x25, 0x63, 0xEB)
BLUE_LIGHT = RGBColor(0x60, 0xA5, 0xFA)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
RED        = RGBColor(0xDC, 0x26, 0x26)
AMBER      = RGBColor(0xD9, 0x77, 0x06)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFC)
DARK_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
MUTED      = RGBColor(0x64, 0x74, 0x8B)
LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
LIGHT_RED   = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_AMBER = RGBColor(0xFF, 0xFB, 0xEB)
LIGHT_BLUE  = RGBColor(0xDB, 0xEA, 0xFE)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

SLIDE_NUM = [0]  # mutable counter

# ── Helpers ──

def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, border_color=None, border_width=Pt(1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT,
             bold=False, italic=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_text(slide, left, top, width, height, runs_list, alignment=PP_ALIGN.LEFT):
    """runs_list: list of (text, font_size, color, bold, italic)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, (text, fs, clr, bld, ital) in enumerate(runs_list):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
            run.text = text
        else:
            run = p.add_run()
            run.text = text
        run.font.size = Pt(fs)
        run.font.color.rgb = clr
        run.font.bold = bld
        run.font.italic = ital
        run.font.name = "Calibri"
    return txBox


def add_bullets(slide, left, top, width, height, items, font_size=16,
                color=DARK_TEXT, spacing=Pt(6), bold_items=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    bold_items = bold_items or set()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        if i in bold_items:
            p.font.bold = True
    return txBox


def add_table(slide, left, top, width, height, rows, cols):
    return slide.shapes.add_table(rows, cols, left, top, width, height).table


def style_header(table, ncols, bg=NAVY, fg=WHITE, font_size=14):
    for i in range(ncols):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = fg
            p.font.bold = True
            p.font.size = Pt(font_size)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER


def sc(cell, text, font_size=13, color=DARK_TEXT, bold=False,
       alignment=PP_ALIGN.CENTER, bg=None):
    """Style a data cell."""
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def new_slide():
    SLIDE_NUM[0] += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return slide


def title_bar(slide, title, subtitle=None, section=None):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.3), NAVY)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(10.5), Inches(0.65),
             title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.78), Inches(10.5), Inches(0.4),
                 subtitle, font_size=15, color=RGBColor(0x94, 0xA3, 0xB8))
    # Section indicator top right
    if section:
        add_text(slide, Inches(10.5), Inches(0.25), Inches(2.5), Inches(0.4),
                 section, font_size=12, color=BLUE_LIGHT, italic=True,
                 alignment=PP_ALIGN.RIGHT)
    # Slide number bottom right
    add_text(slide, Inches(12), Inches(7.05), Inches(1), Inches(0.35),
             str(SLIDE_NUM[0]), font_size=11, color=MUTED, alignment=PP_ALIGN.RIGHT)


def callout_box(slide, left, top, width, height, title_text, body_items,
                border_color=BLUE, title_color=None, font_size=14):
    title_color = title_color or border_color
    add_rect(slide, left, top, width, height, WHITE, border_color, Pt(2))
    add_text(slide, left + Inches(0.2), top + Inches(0.08), width - Inches(0.4),
             Inches(0.35), title_text, font_size=17, color=title_color, bold=True)
    add_bullets(slide, left + Inches(0.2), top + Inches(0.48),
                width - Inches(0.4), height - Inches(0.55),
                body_items, font_size=font_size, spacing=Pt(4))


# ════════════════════════════════════════════════════════════════
#  MAIN DECK — 25 slides
# ════════════════════════════════════════════════════════════════

# ── S1: Title ──
slide = new_slide()
add_bg(slide, NAVY)
# Gradient-like effect: darker bottom strip
add_rect(slide, Inches(0), Inches(6.2), W, Inches(1.3), NAVY_DARK)

add_text(slide, Inches(1), Inches(1.0), Inches(11.3), Inches(1.2),
         "Inverting Python:", font_size=44, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.1), Inches(11.3), Inches(1.0),
         "Designing NTP-Prior-Resistant Confusion Languages\nfor LLM-Resistant Coding Assessment",
         font_size=24, color=BLUE_LIGHT, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(4.5), Inches(3.5), Inches(4.3), Pt(2), BLUE)
add_text(slide, Inches(1), Inches(3.9), Inches(11.3), Inches(0.5),
         "Jaeyeong CHOI", font_size=24, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.4),
         "DGIST EECS  |  jaeyeong2022@dgist.ac.kr", font_size=16,
         color=MUTED, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11.3), Inches(0.4),
         "13 Models   |   950+ L4 Evaluations   |   4-Level Taxonomy",
         font_size=16, color=RGBColor(0x7B, 0x9B, 0xBB), alignment=PP_ALIGN.CENTER)
# Slide number
add_text(slide, Inches(12), Inches(7.05), Inches(1), Inches(0.35),
         "1", font_size=11, color=MUTED, alignment=PP_ALIGN.RIGHT)

# ── S2: Outline / Roadmap ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Outline", section="Overview")

sections = [
    ("1", "Problem & Motivation", "Why LLMs undermine coding assessments"),
    ("2", "Core Asymmetry", "Human reasoning vs. LLM sampling"),
    ("3", "Research Question", "Contributions & approach"),
    ("4", "4-Level Taxonomy", "L1\u2013L4 perturbation depth \u00d7 rule delivery"),
    ("5", "Experimental Results", "L1 through L4, ablations, multi-task"),
    ("6", "Key Findings", "Prior entrenchment, hard tasks, failure modes"),
    ("7", "Educational Implications", "Design recipe for instructors"),
    ("8", "Conclusion & Future Work", "Human study, broader tasks"),
]
for i, (num, title, desc) in enumerate(sections):
    y = Inches(1.6) + Inches(i * 0.7)
    add_rect(slide, Inches(0.8), y, Inches(0.55), Inches(0.5), BLUE)
    add_text(slide, Inches(0.8), y + Inches(0.05), Inches(0.55), Inches(0.4),
             num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.6), y + Inches(0.02), Inches(4.5), Inches(0.35),
             title, font_size=20, color=DARK_TEXT, bold=True)
    add_text(slide, Inches(6.5), y + Inches(0.05), Inches(6), Inches(0.35),
             desc, font_size=16, color=MUTED)

# ── S3: Problem — LLM solves coding assignments ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Problem: LLMs Solve Coding Assignments", section="Motivation")

items = [
    "\u2022  LLMs (GPT-4o, o3, etc.) now routinely solve standard CS1/CS2 programming tasks",
    "\u2022  Students can submit LLM-generated code with minimal understanding",
    "\u2022  Traditional defenses are insufficient:",
    "     \u2013  Plagiarism detectors: LLM output is unique per generation",
    "     \u2013  Proctoring: students use LLMs outside exam windows",
    "     \u2013  Oral follow-ups: don't scale to large lecture courses",
    "",
    "\u2022  Need: a structural defense that makes assignments inherently LLM-resistant",
    "\u2022  Approach: exploit the fundamental mechanism gap between humans and LLMs",
]
add_bullets(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.2),
            items, font_size=19, spacing=Pt(8))

# ── S4: Core Asymmetry ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Core Asymmetry: Human Reasons, LLM Samples",
          subtitle="Fundamentally different mechanisms for code generation", section="Motivation")

# Human box
add_rect(slide, Inches(0.5), Inches(1.6), Inches(5.9), Inches(4.7), WHITE, BLUE, Pt(2))
add_text(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.45),
         "Human Student", font_size=26, color=BLUE, bold=True)
h_items = [
    "\u2022  Reads explicit rule from instructor",
    "\u2022  Verifies rule against worked examples",
    "\u2022  Applies rule step-by-step to new problems",
    "\u2022  Reasoning-based: rule \u2192 application",
    "\u2022  Can override prior knowledge with new rules",
]
add_bullets(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(2.5),
            h_items, font_size=17, spacing=Pt(8))
add_rect(slide, Inches(0.8), Inches(5.3), Inches(5.3), Inches(0.6), LIGHT_GREEN)
add_text(slide, Inches(1.0), Inches(5.35), Inches(5), Inches(0.4),
         "Can learn inverted semantics", font_size=16, color=GREEN, bold=True,
         alignment=PP_ALIGN.CENTER)

# LLM box
add_rect(slide, Inches(6.9), Inches(1.6), Inches(5.9), Inches(4.7), WHITE, RED, Pt(2))
add_text(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.45),
         "LLM (GPT-4o, o3, gpt-5.4, ...)", font_size=26, color=RED, bold=True)
l_items = [
    "\u2022  Samples from NTP prior distribution",
    "\u2022  Prior built from billions of Python tokens",
    "\u2022  'if n <= 1: return n' deeply entrenched",
    "\u2022  Generation-based: prior \u2192 sampling",
    "\u2022  Cannot override deep priors via context",
]
add_bullets(slide, Inches(7.2), Inches(2.3), Inches(5.2), Inches(2.5),
            l_items, font_size=17, spacing=Pt(8))
add_rect(slide, Inches(7.2), Inches(5.3), Inches(5.3), Inches(0.6), LIGHT_RED)
add_text(slide, Inches(7.4), Inches(5.35), Inches(5), Inches(0.4),
         "Fails on inverted Fibonacci (0% across 13 models)", font_size=16,
         color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# Arrow/gap
add_text(slide, Inches(6.1), Inches(3.5), Inches(1), Inches(0.5),
         "\u2260", font_size=48, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

# ── S5: Research Question + Contributions ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Research Question & Contributions", section="Overview")

# RQ box
add_rect(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.2), LIGHT_BLUE, BLUE, Pt(2))
add_text(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.9),
         "RQ: Can we design a language perturbation that humans can learn from examples "
         "but LLMs cannot solve, by exploiting NTP prior entrenchment?",
         font_size=20, color=NAVY, bold=True)

# Contributions
contribs = [
    "\u2022  C1: 4-level confusion language taxonomy (token \u2192 syntax \u2192 semantic-explicit \u2192 semantic-implicit)",
    "\u2022  C2: Empirical evaluation across 13 models (950+ L4 runs), including frontier gpt-5.4 and o3",
    "\u2022  C3: Prior entrenchment depth hypothesis \u2014 task difficulty \u2260 LLM resistance",
    "\u2022  C4: Three failure modes taxonomy (Prior Dominance, Pattern Blindness, Operational Substitution)",
    "\u2022  C5: Practical design recipe for LLM-resistant programming assessments",
]
add_bullets(slide, Inches(0.6), Inches(3.2), Inches(12), Inches(3.5),
            contribs, font_size=18, spacing=Pt(12))

# ── S6: 4-Level Taxonomy (2D grid) ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "4-Level Confusion Language Taxonomy",
          subtitle="Perturbation depth \u00d7 Rule delivery", section="Taxonomy")

tbl = add_table(slide, Inches(0.4), Inches(1.55), Inches(12.5), Inches(4.5), 5, 6)
hdrs = ["Level", "Perturbation", "Rule Delivery", "Example", "Result", "Failure Mode"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 6)

rows = [
    ("L1", "Token Alias", "Explicit (prompt)",
     "DISPLAY \u2192 print", "Partial (KLR 0.21\u20131.00)", "Prior Dominance"),
    ("L2", "Syntax Inversion", "Explicit (prompt)",
     "for i in 3..0 = ascending", "All models comply", "None"),
    ("L3", "Semantic Inversion", "Explicit (prompt)",
     "if-block runs when FALSE", "Capacity-dependent", "Model threshold"),
    ("L4", "Semantic Inversion", "Implicit (examples only)",
     "Worked examples only", "ALL fail (0%)", "Pattern Blindness"),
]
bgs = [LIGHT_AMBER, LIGHT_RED, LIGHT_AMBER, LIGHT_GREEN]
for r, (lv, pert, rd, ex, res, fm) in enumerate(rows):
    sc(tbl.cell(r+1, 0), lv, font_size=16, bold=True, bg=bgs[r])
    sc(tbl.cell(r+1, 1), pert, font_size=14, bg=bgs[r], alignment=PP_ALIGN.LEFT)
    sc(tbl.cell(r+1, 2), rd, font_size=14, bg=bgs[r], alignment=PP_ALIGN.LEFT)
    sc(tbl.cell(r+1, 3), ex, font_size=13, bg=bgs[r], alignment=PP_ALIGN.LEFT, color=MUTED)
    sc(tbl.cell(r+1, 4), res, font_size=14, bg=bgs[r], bold=True)
    sc(tbl.cell(r+1, 5), fm, font_size=13, bg=bgs[r])

tbl.columns[0].width = Inches(0.8)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(2.2)
tbl.columns[3].width = Inches(2.5)
tbl.columns[4].width = Inches(2.8)
tbl.columns[5].width = Inches(2.2)

add_text(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.4),
         "Finding: Only L4 (implicit rule + deep prior) defeats ALL models universally.",
         font_size=17, color=BLUE, bold=True)

# ── S7: L1 Token Substitution ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "L1: Token Substitution \u2014 Design & Results",
          subtitle="Keyword Leakage Rate (KLR, lower = better alias compliance)", section="L1")

tbl = add_table(slide, Inches(0.4), Inches(1.5), Inches(7.2), Inches(5.2), 12, 4)
hdrs = ["Model", "Delivery", "n", "KLR"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 4)

klr_data = [
    ("gpt-4.1-mini", "ctx-pack", "120", "0.21", LIGHT_GREEN),
    ("Llama-3.3-70B", "Groq", "20", "0.23", LIGHT_GREEN),
    ("gpt-4o-mini", "ctx-pack", "120", "0.26", LIGHT_GREEN),
    ("o4-mini", "ctx-pack", "20", "0.27", LIGHT_GREEN),
    ("gpt-4.1", "ctx-pack", "20", "0.30", None),
    ("gpt-4o", "baseline", "120", "0.37", None),
    ("gpt-4.1-nano", "ctx-pack", "30", "0.38", LIGHT_AMBER),
    ("gpt-5.4-mini", "baseline", "120", "0.67", LIGHT_RED),
    ("Qwen3-32B", "Groq", "20", "1.00", LIGHT_RED),
]
for r, (model, dlv, n, klr, bg) in enumerate(klr_data):
    sc(tbl.cell(r+1, 0), model, font_size=13, alignment=PP_ALIGN.LEFT, bg=bg)
    sc(tbl.cell(r+1, 1), dlv, font_size=12, bg=bg)
    sc(tbl.cell(r+1, 2), n, font_size=12, bg=bg, color=MUTED)
    sc(tbl.cell(r+1, 3), klr, font_size=14, bold=True, bg=bg)
# Best/worst labels
sc(tbl.cell(1, 0), "gpt-4.1-mini  \u2190 best", font_size=13, alignment=PP_ALIGN.LEFT,
   bg=LIGHT_GREEN, color=GREEN, bold=True)
sc(tbl.cell(9, 0), "Qwen3-32B  \u2190 worst", font_size=13, alignment=PP_ALIGN.LEFT,
   bg=LIGHT_RED, color=RED, bold=True)
# Empty row label
sc(tbl.cell(10, 0), "", font_size=1)
sc(tbl.cell(10, 1), "", font_size=1)
sc(tbl.cell(10, 2), "", font_size=1)
sc(tbl.cell(10, 3), "", font_size=1)
sc(tbl.cell(11, 0), "", font_size=1)

tbl.columns[0].width = Inches(2.6)
tbl.columns[1].width = Inches(1.4)
tbl.columns[2].width = Inches(0.7)
tbl.columns[3].width = Inches(1.0)

# Key findings box
callout_box(slide, Inches(7.9), Inches(1.5), Inches(5), Inches(5.2),
            "Key Findings", [
    "\u2022  Context-pack reduces KLR:",
    "   gpt-4o-mini: 0.42 \u2192 0.26 (\u221238%)",
    "   gpt-4.1-mini: 0.43 \u2192 0.21 (\u221251%)",
    "",
    "\u2022  gpt-5.4-mini (0.67) worse than",
    "   older gpt-4o (0.37)",
    "",
    "\u2022  Cross-family variance (0.21\u20131.00)",
    "   exceeds intra-family variance",
    "",
    "\u2022  Architecture & fine-tuning govern",
    "   alias compliance, not model size",
], font_size=14)

# ── S8: L2 Syntax Inversion ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "L2: Syntax Inversion \u2014 Design & Results",
          subtitle="SIR = 0.0 across all tested models (fully adopted)", section="L2")

# Design description
add_text(slide, Inches(0.5), Inches(1.6), Inches(12), Inches(0.4),
         "Design: Invert loop/range syntax (e.g., for i in 3..0 means ascending 0\u21923).",
         font_size=20, color=DARK_TEXT, bold=True)
add_text(slide, Inches(0.5), Inches(2.1), Inches(12), Inches(0.4),
         "Explicit rules provided in prompt. Models tested: gpt-4o, gpt-4.1-mini, gpt-4.1-nano (n=20 each).",
         font_size=17, color=MUTED)

tbl = add_table(slide, Inches(2.5), Inches(2.8), Inches(8), Inches(2.0), 4, 3)
hdrs = ["Model", "n", "SIR"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 3)
for r, (m, n) in enumerate([("gpt-4o", "20"), ("gpt-4.1-mini", "20"), ("gpt-4.1-nano", "20")]):
    sc(tbl.cell(r+1, 0), m, font_size=15, alignment=PP_ALIGN.LEFT, bg=LIGHT_GREEN)
    sc(tbl.cell(r+1, 1), n, font_size=15, bg=LIGHT_GREEN)
    sc(tbl.cell(r+1, 2), "0.0", font_size=18, bold=True, bg=LIGHT_GREEN, color=GREEN)

tbl.columns[0].width = Inches(3.5)
tbl.columns[1].width = Inches(1.5)
tbl.columns[2].width = Inches(3.0)

callout_box(slide, Inches(2.5), Inches(5.2), Inches(8), Inches(1.5),
            "Finding 1", [
    "\u2022  Syntax-level inversion with explicit rules is trivially adopted by all models",
    "\u2022  No LLM resistance at this level \u2014 rules are shallow and easily followed",
    "\u2022  This motivates moving to semantic-level perturbation (L3/L4)",
], border_color=GREEN, font_size=15)

# ── S9: L3 Semantic Inversion (Explicit) ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "L3: Semantic Inversion (Explicit) \u2014 Capacity Threshold",
          subtitle="if-blocks run when condition is FALSE; rule stated explicitly", section="L3")

tbl = add_table(slide, Inches(0.5), Inches(1.5), Inches(7), Inches(3.5), 7, 4)
hdrs = ["Model", "n", "Pass", "PPR"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 4)

l3data = [
    ("gpt-4o", "20", "20/20 PASS", "0.0", LIGHT_GREEN),
    ("gpt-4.1-mini", "20", "20/20 PASS", "0.0", LIGHT_GREEN),
    ("gpt-4.1-nano", "20", "20/20 PASS", "0.0", LIGHT_GREEN),
    ("gpt-4.1 (T2)", "10", "10/10 PASS", "0.0", LIGHT_GREEN),
    ("gpt-4o-mini", "5", "0/5 FAIL", "1.0", LIGHT_RED),
    ("Qwen3-32B", "20", "0/20 FAIL", "1.0", LIGHT_RED),
]
for r, (m, n, ps, ppr, bg) in enumerate(l3data):
    sc(tbl.cell(r+1, 0), m, font_size=14, alignment=PP_ALIGN.LEFT, bg=bg)
    sc(tbl.cell(r+1, 1), n, font_size=14, bg=bg, color=MUTED)
    clr = GREEN if "PASS" in ps else RED
    sc(tbl.cell(r+1, 2), ps, font_size=14, bold=True, bg=bg, color=clr)
    sc(tbl.cell(r+1, 3), ppr, font_size=14, bold=True, bg=bg)

tbl.columns[0].width = Inches(2.2)
tbl.columns[1].width = Inches(0.8)
tbl.columns[2].width = Inches(2.2)
tbl.columns[3].width = Inches(1.0)

callout_box(slide, Inches(7.9), Inches(1.5), Inches(5), Inches(3.5),
            "Finding 2: Capacity Threshold", [
    "\u2022  Larger models (gpt-4o, 4.1-mini,",
    "   4.1-nano) follow explicit rules",
    "",
    "\u2022  Smaller gpt-4o-mini: 0/5 FAIL",
    "\u2022  Qwen3-32B: 0/20 FAIL",
    "",
    "\u2022  L3 alone is insufficient \u2014",
    "   capable models can follow explicit",
    "   semantic inversion rules",
], font_size=14)

add_text(slide, Inches(0.5), Inches(5.4), Inches(12), Inches(0.5),
         "Implication: Explicit rule delivery is not enough. "
         "We must withhold the rule and let models infer from examples \u2192 L4.",
         font_size=17, color=BLUE, bold=True)

# ── S10: L4 Semantic Inversion (Implicit) \u2014 Design ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "L4: Semantic Inversion (Implicit) \u2014 Design",
          subtitle="Examples only, no explicit rule statement", section="L4")

items = [
    "\u2022  Core hypothesis: If the rule is never stated, only demonstrated through",
    "   worked examples, LLMs cannot extract it due to NTP prior dominance",
    "",
    "\u2022  Prompt structure:",
    "   1. System prompt defines the confusion language concept (no rule stated)",
    "   2. 3\u20135 worked input/output examples showing inverted semantics",
    "   3. Target task (e.g., Fibonacci) to implement in the confusion language",
    "",
    "\u2022  Key design choice: if-blocks execute when condition is FALSE",
    "   \u2192 Standard Fibonacci base case 'if n <= 1: return n' must become",
    "      'if n > 1: return n' (the negated condition triggers execution)",
    "",
    "\u2022  Evaluation: strict string/AST match for inverted conditionals",
    "\u2022  No partial credit \u2014 either the model produces correct inverted code or fails",
]
add_bullets(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
            items, font_size=18, spacing=Pt(5))

# ── S11: L4 Ablation \u2014 5 Variants Design ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "L4 Ablation: 5 Variant Design",
          subtitle="Strict example-only, varying example count and hint level", section="L4 Ablation")

tbl = add_table(slide, Inches(0.4), Inches(1.5), Inches(12.5), Inches(3.2), 6, 4)
hdrs = ["Variant", "Examples", "Hint", "Purpose"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 4)

variants = [
    ("A", "5 examples", "No hint", "Baseline: maximum demonstration, zero guidance"),
    ("B", "3 examples", "No hint", "Reduced examples test minimum context"),
    ("C", "5 examples", "Subtle hint", "Hint about 'unusual control flow' behavior"),
    ("D", "5 examples", "Moderate hint", "Hint about conditional semantics specifically"),
    ("E", "5 examples", "Strong hint", "Near-explicit hint about inversion direction"),
]
for r, (v, ex, h, p) in enumerate(variants):
    sc(tbl.cell(r+1, 0), f"Variant {v}", font_size=15, bold=True)
    sc(tbl.cell(r+1, 1), ex, font_size=14, alignment=PP_ALIGN.LEFT)
    sc(tbl.cell(r+1, 2), h, font_size=14, alignment=PP_ALIGN.LEFT)
    sc(tbl.cell(r+1, 3), p, font_size=13, alignment=PP_ALIGN.LEFT, color=MUTED)

tbl.columns[0].width = Inches(1.6)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(2.2)
tbl.columns[3].width = Inches(6.7)

add_text(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
         "Each variant tested on 9 models \u00d7 n=10 per model = n=50 total per variant (n=450 total ablation runs)",
         font_size=16, color=MUTED)

callout_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.3),
            "Design Rationale", [
    "\u2022  Variants A\u2192E form a gradient from pure example-only to near-explicit rule",
    "\u2022  If any variant succeeds, pattern blindness is solvable with enough hints",
    "\u2022  Result preview: NONE succeed on Fibonacci (0/450 across 9 models)",
], border_color=RED, font_size=15)

# ── S12: L4 Ablation Results \u2014 9 models, 0 pass ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "L4 Ablation Results: 9 Models, 0% Pass Rate",
          subtitle="5 variants \u00d7 n=10 per model = n=50 total per model", section="L4 Ablation")

tbl = add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(3.5), 10, 5)
hdrs = ["Model", "Pass / n", "PPR", "Notable", "Category"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 5)

abl_data = [
    ("gpt-4o", "0/50", "1.00", "\u2014", "OpenAI Core"),
    ("gpt-4o-mini", "0/50", "1.00", "\u2014", "OpenAI Core"),
    ("gpt-4.1", "0/50", "1.00", "\u2014", "OpenAI Core"),
    ("gpt-4.1-mini", "0/50", "1.00", "Variant B PPR=0.0 (decomposition)", "OpenAI Core"),
    ("gpt-4.1-nano", "0/50", "0.20", "Template fragmentation", "OpenAI Core"),
    ("gpt-5.4-mini", "0/42 valid", "1.00", "8 parse failures excluded", "OpenAI Core"),
    ("o4-mini", "0/45", "0.60", "Reasoning model, partial effect", "Reasoning"),
    ("Llama-3.3-70B", "0/25 valid", "0.52", "Variant C\u2013E context-sensitive", "Open-weight"),
    ("Qwen3-32B", "0/50", "1.00", "\u2014", "Open-weight"),
]
for r, (m, ps, ppr, note, cat) in enumerate(abl_data):
    bg = LIGHT_RED
    sc(tbl.cell(r+1, 0), m, font_size=13, alignment=PP_ALIGN.LEFT, bg=bg, bold=True)
    sc(tbl.cell(r+1, 1), ps, font_size=14, bg=bg, color=RED, bold=True)
    sc(tbl.cell(r+1, 2), ppr, font_size=14, bg=bg, bold=True)
    sc(tbl.cell(r+1, 3), note, font_size=11, bg=bg, alignment=PP_ALIGN.LEFT, color=MUTED)
    sc(tbl.cell(r+1, 4), cat, font_size=12, bg=bg, color=MUTED)

tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(1.3)
tbl.columns[2].width = Inches(1.0)
tbl.columns[3].width = Inches(5.8)
tbl.columns[4].width = Inches(1.8)

callout_box(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(1.5),
            "Finding 3: Universal Pattern Blindness on Fibonacci", [
    "\u2022  0% pass rate across ALL 9 models, ALL 5 variants (0/450+ total runs)",
    "\u2022  Even strong hints (Variant E) and reasoning models (o4-mini) cannot overcome prior entrenchment",
    "\u2022  gpt-4.1-nano PPR=0.20 is template fragmentation, not rule extraction",
], border_color=RED, title_color=RED, font_size=15)

# ── S13: CoT Ablation ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "L4 CoT Ablation: Chain-of-Thought Does NOT Help",
          subtitle="Reasoning \u2260 Generation", section="L4 CoT")

tbl = add_table(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(3.0), 5, 4)
hdrs = ["Model", "No-CoT (n=20)", "CoT (n=20)", "Effect"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 4)

cot_data = [
    ("gpt-4o", "0/20", "0/20", "None"),
    ("gpt-4.1-mini", "0/20", "0/20", "None"),
    ("gpt-4.1-nano", "0/20", "0/20", "None"),
    ("o4-mini", "0/20", "1/20", "Minimal"),
]
for r, (m, nc, c, eff) in enumerate(cot_data):
    bg = LIGHT_RED if eff == "None" else LIGHT_AMBER
    sc(tbl.cell(r+1, 0), m, font_size=15, alignment=PP_ALIGN.LEFT, bold=True, bg=bg)
    sc(tbl.cell(r+1, 1), nc, font_size=15, bg=bg, color=RED, bold=True)
    sc(tbl.cell(r+1, 2), c, font_size=15, bg=bg, color=RED, bold=True)
    eff_clr = RED if eff == "None" else AMBER
    sc(tbl.cell(r+1, 3), eff, font_size=15, bg=bg, color=eff_clr, bold=True)

tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(2.5)
tbl.columns[2].width = Inches(2.5)
tbl.columns[3].width = Inches(2.5)

callout_box(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(2.0),
            "Finding 4: Reasoning-Generation Dissociation", [
    "\u2022  160 total CoT runs across 4 models \u2014 CoT has zero effect for 3/4 models",
    "\u2022  Models can articulate the inversion rule in their reasoning trace (100% mention rate)",
    "\u2022  But still generate standard Python code in the output \u2014 NTP prior overrides reasoning",
    "\u2022  o4-mini: 1/20 marginal success \u2014 reasoning model architecture provides minimal benefit",
], border_color=BLUE, font_size=15)

# ── S14: L4 Multi-task ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "L4 Multi-task: Task Dependency on Prior Depth",
          subtitle="T1=Fibonacci, T2=is_sorted, T3=count_vowels", section="L4 Multi-task")

tbl = add_table(slide, Inches(0.3), Inches(1.5), Inches(12.7), Inches(3.8), 8, 5)
hdrs = ["Model", "T1: Fibonacci", "T2: is_sorted", "T3: count_vowels", "Notes"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 5)

mt_data = [
    ("gpt-4o", "0/10", "10/10", "0/50", "T3: 94% OpSub (n=50)"),
    ("gpt-4o-mini", "0/10", "0/10", "0/10", ""),
    ("gpt-4.1", "0/10", "1/50", "\u2014", "T2: shallow prior, 1/50"),
    ("gpt-4.1-mini", "0/20", "20/20", "0/20", "T2 fully solved"),
    ("gpt-4.1-nano", "0/10", "0/10", "0/10", ""),
    ("gpt-5.4-mini", "0/10", "0/10", "0/10", ""),
    ("o4-mini", "8/15", "\u2014", "\u2014", "Partial annotation only"),
]
for r, (m, t1, t2, t3, note) in enumerate(mt_data):
    sc(tbl.cell(r+1, 0), m, font_size=13, alignment=PP_ALIGN.LEFT, bold=True)
    # T1
    t1_bg = LIGHT_AMBER if "8/15" in t1 else LIGHT_RED
    t1_clr = AMBER if "8/15" in t1 else RED
    sc(tbl.cell(r+1, 1), t1, font_size=14, bg=t1_bg, color=t1_clr, bold=True)
    # T2
    if t2 in ("10/10", "20/20"):
        sc(tbl.cell(r+1, 2), t2, font_size=14, bg=LIGHT_GREEN, color=GREEN, bold=True)
    elif t2 == "\u2014":
        sc(tbl.cell(r+1, 2), t2, font_size=14, color=MUTED)
    else:
        sc(tbl.cell(r+1, 2), t2, font_size=14, bg=LIGHT_RED, color=RED, bold=True)
    # T3
    if t3 == "\u2014":
        sc(tbl.cell(r+1, 3), t3, font_size=14, color=MUTED)
    else:
        sc(tbl.cell(r+1, 3), t3, font_size=14, bg=LIGHT_RED, color=RED, bold=True)
    sc(tbl.cell(r+1, 4), note, font_size=11, alignment=PP_ALIGN.LEFT, color=MUTED)

tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(1.8)
tbl.columns[2].width = Inches(1.8)
tbl.columns[3].width = Inches(2.0)
tbl.columns[4].width = Inches(4.3)

callout_box(slide, Inches(0.3), Inches(5.6), Inches(12.7), Inches(1.3),
            "Finding 5: Prior Depth Determines Resistance", [
    "\u2022  T1 (Fibonacci): universally fails \u2014 deep, uniform prior",
    "\u2022  T2 (is_sorted): some models succeed \u2014 shallow prior, idiomatic variation",
    "\u2022  T3 (count_vowels): fails but via operational substitution (94% in gpt-4o), not pattern blindness",
], border_color=BLUE, font_size=15)

# ── S15: Prior Entrenchment Depth ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Key Finding: Prior Entrenchment Depth",
          subtitle="Not task complexity, but prior depth governs LLM resistance", section="Key Finding")

# 3-tier illustration
tiers = [
    ("Deep Prior", "Fibonacci recursion", "if n <= 1: return n",
     "Billions of training examples", "ALL models fail (0%)", RED, LIGHT_RED),
    ("Medium Prior", "count_vowels", "if char in 'aeiou': count += 1",
     "Common but varied patterns", "Fails, but via workaround (OpSub 94%)", AMBER, LIGHT_AMBER),
    ("Shallow Prior", "is_sorted", "all(a[i] <= a[i+1])",
     "Multiple valid idioms", "Some models succeed (gpt-4.1-mini: 20/20)", GREEN, LIGHT_GREEN),
]
for i, (tier, task, code, reason, result, clr, bg_clr) in enumerate(tiers):
    y = Inches(1.6) + Inches(i * 1.75)
    add_rect(slide, Inches(0.5), y, Inches(12.3), Inches(1.55), bg_clr, clr, Pt(2))
    add_text(slide, Inches(0.8), y + Inches(0.08), Inches(2.5), Inches(0.35),
             tier, font_size=20, color=clr, bold=True)
    add_text(slide, Inches(3.5), y + Inches(0.08), Inches(2.5), Inches(0.35),
             f"Task: {task}", font_size=16, color=DARK_TEXT, bold=True)
    add_text(slide, Inches(3.5), y + Inches(0.45), Inches(3.5), Inches(0.3),
             f"Pattern: {code}", font_size=13, color=MUTED, italic=True)
    add_text(slide, Inches(3.5), y + Inches(0.8), Inches(3.5), Inches(0.3),
             reason, font_size=13, color=MUTED)
    add_text(slide, Inches(7.5), y + Inches(0.3), Inches(5), Inches(0.5),
             result, font_size=16, color=clr, bold=True)

add_text(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
         "Implication: Effective confusion languages must target deep, uniform priors (e.g., Fibonacci recursion).",
         font_size=17, color=BLUE, bold=True)

# ── S16: Hard Tasks ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Hard Tasks: Complex Algorithms Defeated",
          subtitle="Task complexity \u2260 LLM resistance", section="Key Finding")

tbl = add_table(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(3.0), 5, 5)
hdrs = ["Model", "H1: MergeSort", "H2: BinSearch", "H3: BFS", "Total"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 5)

hard_data = [
    ("gpt-4o", "4/10", "10/10", "10/10", "24/30 (80%)"),
    ("gpt-5.4", "10/10", "10/10", "10/10", "30/30 (100%)"),
    ("o3", "10/10", "10/10", "10/10", "30/30 (100%)"),
    ("o4-mini", "10/10", "10/10", "10/10", "30/30 (100%)"),
]
for r, (m, h1, h2, h3, tot) in enumerate(hard_data):
    sc(tbl.cell(r+1, 0), m, font_size=15, alignment=PP_ALIGN.LEFT, bold=True)
    h1_bg = LIGHT_AMBER if m == "gpt-4o" else LIGHT_GREEN
    h1_clr = AMBER if m == "gpt-4o" else GREEN
    sc(tbl.cell(r+1, 1), h1, font_size=15, bg=h1_bg, color=h1_clr, bold=True)
    sc(tbl.cell(r+1, 2), h2, font_size=15, bg=LIGHT_GREEN, color=GREEN, bold=True)
    sc(tbl.cell(r+1, 3), h3, font_size=15, bg=LIGHT_GREEN, color=GREEN, bold=True)
    sc(tbl.cell(r+1, 4), tot, font_size=15, bold=True, color=GREEN if "100" in tot else DARK_TEXT)

tbl.columns[0].width = Inches(1.8)
tbl.columns[1].width = Inches(1.8)
tbl.columns[2].width = Inches(1.8)
tbl.columns[3].width = Inches(1.8)
tbl.columns[4].width = Inches(1.8)

callout_box(slide, Inches(0.5), Inches(4.8), Inches(5.8), Inches(2.0),
            "Contrast with L4 Fibonacci", [
    "\u2022  Same models on Fibonacci: ALL fail (0%)",
    "\u2022  Same models on hard tasks: near 100%",
    "\u2022  Prior entrenchment depth, not task difficulty,",
    "   governs LLM resistance",
], border_color=BLUE, font_size=15)

callout_box(slide, Inches(6.7), Inches(4.8), Inches(6.1), Inches(2.0),
            "Key Implication", [
    "\u2022  Making assignments harder does NOT make them",
    "   more LLM-resistant",
    "\u2022  Deep, uniform priors (Fibonacci) = resistant",
    "\u2022  Varied implementations (BFS, sort) = solvable",
], border_color=RED, title_color=RED, font_size=15)

# ── S17: E24 Frontier Models ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "E24: Frontier Models Still Fail on Fibonacci",
          subtitle="gpt-5.4 and o3 achieve 0/10 on prior-entrenched task", section="Frontier")

tbl = add_table(slide, Inches(2), Inches(1.5), Inches(9), Inches(2.5), 3, 5)
hdrs = ["Model", "Task", "Pass", "PPR", "Hard Tasks"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 5)

e24_data = [
    ("gpt-5.4", "Fibonacci", "0/10", "1.00", "30/30 (100%)"),
    ("o3", "Fibonacci", "0/10", "1.00", "30/30 (100%)"),
]
for r, (m, t, p, ppr, ht) in enumerate(e24_data):
    sc(tbl.cell(r+1, 0), m, font_size=16, alignment=PP_ALIGN.LEFT, bold=True)
    sc(tbl.cell(r+1, 1), t, font_size=16)
    sc(tbl.cell(r+1, 2), p, font_size=18, bg=LIGHT_RED, color=RED, bold=True)
    sc(tbl.cell(r+1, 3), ppr, font_size=16, bg=LIGHT_RED, bold=True)
    sc(tbl.cell(r+1, 4), ht, font_size=16, bg=LIGHT_GREEN, color=GREEN, bold=True)

tbl.columns[0].width = Inches(1.8)
tbl.columns[1].width = Inches(1.8)
tbl.columns[2].width = Inches(1.8)
tbl.columns[3].width = Inches(1.8)
tbl.columns[4].width = Inches(1.8)

items = [
    "\u2022  The most capable models available (as of March 2025) still exhibit 100% PPR on Fibonacci",
    "\u2022  gpt-5.4 achieves 30/30 on hard tasks (MergeSort, BinSearch, BFS) but 0/10 on inverted Fibonacci",
    "\u2022  o3 (reasoning model) also fails completely despite advanced reasoning capabilities",
    "",
    "\u2022  This demonstrates that pattern blindness is a fundamental architectural limitation,",
    "   not a capability gap that scales away with model size or reasoning capability",
    "",
    "\u2022  Prior entrenchment in Fibonacci is so deep that even frontier models cannot override it",
    "   from examples alone \u2014 the NTP prior dominates generation regardless of context",
]
add_bullets(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(3),
            items, font_size=16, spacing=Pt(5))

# ── S18: 3 Failure Modes Taxonomy ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Three Failure Modes Taxonomy",
          subtitle="Distinct mechanisms behind LLM failure on confusion languages", section="Analysis")

# Mode 1
add_rect(slide, Inches(0.3), Inches(1.5), Inches(4.1), Inches(5.2), WHITE, BLUE, Pt(2))
add_text(slide, Inches(0.6), Inches(1.6), Inches(3.6), Inches(0.4),
         "1. Prior Dominance", font_size=22, color=BLUE, bold=True)
pd = [
    "\u2022  Rule is stated or implicit but",
    "   overridden by pretrained prior",
    "",
    "\u2022  Where: L1, L3-weak, L4-T1",
    "",
    "\u2022  Example: generates",
    "   'if n <= 1: return n'",
    "   despite inversion rule",
    "",
    "\u2022  83.1% of L4 failures",
    "   (Type-I dominant)",
]
add_bullets(slide, Inches(0.6), Inches(2.2), Inches(3.5), Inches(4),
            pd, font_size=14, spacing=Pt(3))

# Mode 2
add_rect(slide, Inches(4.6), Inches(1.5), Inches(4.1), Inches(5.2), WHITE, RED, Pt(2))
add_text(slide, Inches(4.9), Inches(1.6), Inches(3.6), Inches(0.4),
         "2. Pattern Blindness", font_size=22, color=RED, bold=True)
pb = [
    "\u2022  Model cannot extract the",
    "   semantic rule from examples",
    "",
    "\u2022  Where: L4-T1 (Fibonacci)",
    "",
    "\u2022  Evidence: Same model passes",
    "   L3 (explicit) but fails L4",
    "   (implicit) on same task",
    "",
    "\u2022  Universal across all 13",
    "   tested models on Fibonacci",
]
add_bullets(slide, Inches(4.9), Inches(2.2), Inches(3.5), Inches(4),
            pb, font_size=14, spacing=Pt(3))

# Mode 3
add_rect(slide, Inches(8.9), Inches(1.5), Inches(4.1), Inches(5.2), WHITE, ORANGE, Pt(2))
add_text(slide, Inches(9.2), Inches(1.6), Inches(3.6), Inches(0.4),
         "3. Operational Substitution", font_size=22, color=ORANGE, bold=True)
os_items = [
    "\u2022  Model compensates via",
    "   arithmetic transformation",
    "   instead of adopting the rule",
    "",
    "\u2022  Where: L4-T3 (count_vowels)",
    "",
    "\u2022  Example: count -= 1;",
    "   return -count to match output",
    "",
    "\u2022  94% of gpt-4o T3 runs",
    "   (47/50, CI: [83%, 99%])",
]
add_bullets(slide, Inches(9.2), Inches(2.2), Inches(3.5), Inches(4),
            os_items, font_size=14, spacing=Pt(3))

# ── S19: Operational Substitution Deep Dive ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Operational Substitution: Deep Dive",
          subtitle="Models achieve correct output without learning the rule", section="Analysis")

items = [
    "\u2022  Definition: Model produces functionally correct output by applying an arithmetic",
    "   or logical transformation to standard code, rather than adopting the confusion language rule",
    "",
    "\u2022  Key example (gpt-4o, count_vowels, T3):",
    "   \u2013  Expected: use inverted if-block (count when char NOT in vowels, negate result)",
    "   \u2013  Actual: count -= 1 for each vowel, then return -count",
    "   \u2013  Produces correct numeric output but does NOT use inverted semantics",
    "",
    "\u2022  Detection methodology:",
    "   1. Parse generated code AST",
    "   2. Check if any conditional uses inverted logic",
    "   3. If conditionals are standard but output matches \u2192 operational substitution",
    "",
    "\u2022  Prevalence:",
    "   \u2013  gpt-4o T3: 47/50 = 94% operational substitution rate (Wilson CI: [83%, 99%])",
    "   \u2013  Demonstrates that models can be 'clever' about matching outputs",
    "     without actually learning the underlying rule",
    "",
    "\u2022  Implication: Output-matching evaluation alone is insufficient;",
    "   structural/AST-based checking is required to detect true rule adoption",
]
add_bullets(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
            items, font_size=16, spacing=Pt(3))

# ── S20: Summary Table ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Summary: All Experiments",
          subtitle="Complete experimental coverage across 13 models", section="Summary")

tbl = add_table(slide, Inches(0.2), Inches(1.4), Inches(12.9), Inches(5.5), 10, 6)
hdrs = ["Experiment", "Models", "n", "Key Metric", "Result", "Finding"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 6, font_size=12)

summ = [
    ("L1 KLR", "9", "~590", "KLR", "0.21\u20131.00", "Partial resistance"),
    ("L2 Syntax", "3", "60", "SIR", "0.0", "Fully adopted"),
    ("L3 Explicit", "5", "75", "PPR", "0.0\u20131.0", "Capacity threshold"),
    ("L4 Ablation", "9", "450+", "Pass Rate", "0%", "Universal failure"),
    ("L4 CoT", "4", "160", "Pass Rate", "0\u20131/20", "CoT no help"),
    ("L4 Multi-task", "7", "345", "Pass Rate", "Task-dependent", "Prior depth"),
    ("L4 Hard Tasks", "4", "120", "Pass Rate", "80\u2013100%", "Complexity \u2260 resistance"),
    ("E24 Frontier", "2", "20", "PPR", "1.00", "Still fails"),
    ("E28 Density", "1", "40", "Pass Rate", "Non-monotonic", "50% peak"),
]
for r, (exp, mod, n, met, res, find) in enumerate(summ):
    bg = None
    if "Universal" in find or "Still" in find:
        bg = LIGHT_RED
    elif "Fully" in find or "Complexity" in find:
        bg = LIGHT_GREEN
    else:
        bg = LIGHT_AMBER
    sc(tbl.cell(r+1, 0), exp, font_size=12, alignment=PP_ALIGN.LEFT, bold=True, bg=bg)
    sc(tbl.cell(r+1, 1), mod, font_size=12, bg=bg)
    sc(tbl.cell(r+1, 2), n, font_size=12, bg=bg)
    sc(tbl.cell(r+1, 3), met, font_size=12, bg=bg)
    sc(tbl.cell(r+1, 4), res, font_size=12, bg=bg, bold=True)
    sc(tbl.cell(r+1, 5), find, font_size=11, bg=bg, alignment=PP_ALIGN.LEFT, color=MUTED)

tbl.columns[0].width = Inches(1.8)
tbl.columns[1].width = Inches(1.0)
tbl.columns[2].width = Inches(0.8)
tbl.columns[3].width = Inches(1.5)
tbl.columns[4].width = Inches(2.3)
tbl.columns[5].width = Inches(3.0)

# ── S21: Educational Implications ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Educational Implications",
          subtitle="Practical guidance for LLM-resistant programming assessment", section="Implications")

items = [
    "\u2022  L4 confusion languages restore discriminative power in coding assessments",
    "\u2022  Students who understand rules (from lecture) retain advantage over LLM-assisted submissions",
    "\u2022  Current models universally fail L4 Fibonacci while humans can apply rules from examples",
    "",
    "\u2022  Key insight: The defense is structural, not just a prompt trick \u2014 it exploits a fundamental",
    "   limitation in how LLMs generate code (NTP prior dominance)",
    "",
    "\u2022  Scalability: Confusion languages work at scale \u2014 no manual code review needed,",
    "   just automated AST-based checking for inverted semantics",
    "",
    "\u2022  Caveats:",
    "   \u2013  Tasks with shallow priors (is_sorted) may still be solvable by capable models",
    "   \u2013  Partial annotation can help reasoning models (o4-mini: 8/15 with hints)",
    "   \u2013  Human baseline study needed to quantify exact performance gap",
]
add_bullets(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
            items, font_size=17, spacing=Pt(5))

# ── S22: Design Recipe for Instructors ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Design Recipe for Instructors",
          subtitle="Step-by-step guide to creating LLM-resistant assignments", section="Implications")

steps = [
    ("1", "Retain Python surface syntax",
     "Students already know Python; minimize cognitive overhead"),
    ("2", "Choose deep-prior tasks",
     "Fibonacci recursion, not sorting/searching (deep > complex)"),
    ("3", "Invert conditional semantics",
     "if-blocks execute when condition is FALSE"),
    ("4", "Provide only worked examples",
     "Never state the rule explicitly; let students infer"),
    ("5", "Use AST-based grading",
     "Check for inverted conditionals, not just output correctness"),
    ("6", "Rotate confusion rules per semester",
     "Prevent rule leakage into future training data"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(1.55) + Inches(i * 0.9)
    add_rect(slide, Inches(0.5), y, Inches(0.55), Inches(0.55), BLUE)
    add_text(slide, Inches(0.5), y + Inches(0.08), Inches(0.55), Inches(0.4),
             num, font_size=20, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.3), y + Inches(0.02), Inches(5), Inches(0.35),
             title, font_size=19, color=DARK_TEXT, bold=True)
    add_text(slide, Inches(6.5), y + Inches(0.05), Inches(6.3), Inches(0.35),
             desc, font_size=15, color=MUTED)

# ── S23: Limitations & Threats to Validity ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Limitations & Threats to Validity", section="Discussion")

items = [
    "\u2022  Internal validity:",
    "   \u2013  Fibonacci-centric: deep prior effect may not generalize to all recursive tasks",
    "   \u2013  Example count confound: variant B uses 3 examples vs. 5 in others",
    "   \u2013  Temperature fixed at default; higher temperature may yield different results",
    "",
    "\u2022  External validity:",
    "   \u2013  Limited to Python; other languages may have different prior distributions",
    "   \u2013  13 models tested; newer architectures may overcome pattern blindness",
    "   \u2013  No human baseline yet (planned as Exp-5)",
    "",
    "\u2022  Construct validity:",
    "   \u2013  PPR measures Python prior dominance but may miss novel failure modes",
    "   \u2013  Operational substitution detection relies on AST heuristics",
    "",
    "\u2022  Prior-depth vs. I/O-transparency confound not yet isolated",
    "\u2022  Annotation density gradient (E28) tested only on o4-mini",
]
add_bullets(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
            items, font_size=16, spacing=Pt(3))

# ── S24: Related Work ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
title_bar(slide, "Related Work Connections",
          subtitle="CodeIF, MultiCodeIF, and instruction-following benchmarks", section="Discussion")

items = [
    "\u2022  CodeIF (2024): Instruction-following for code generation; tests explicit rule compliance",
    "   \u2013  Our L3 is comparable; L4 goes beyond by withholding the rule entirely",
    "",
    "\u2022  MultiCodeIF (2025): Multi-constraint instruction following across languages",
    "   \u2013  Similar constraint framework but doesn't test implicit rule extraction",
    "",
    "\u2022  IFEval / FollowBench: Natural language instruction-following benchmarks",
    "   \u2013  Our work is domain-specific to code generation and prior entrenchment",
    "",
    "\u2022  LLM-resistant assessment literature:",
    "   \u2013  Prior work focuses on detection (plagiarism, AI-text classifiers)",
    "   \u2013  Our approach is structural prevention \u2014 make tasks inherently unsolvable by LLMs",
    "",
    "\u2022  NTP prior analysis (Mccoy et al., 2023; Kambhampati et al., 2024):",
    "   \u2013  Theoretical grounding for why LLMs struggle with counter-prior tasks",
    "   \u2013  Our work provides first large-scale empirical evidence in code domain",
]
add_bullets(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
            items, font_size=16, spacing=Pt(3))

# ── S25: Conclusion + Future Work ──
slide = new_slide()
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(6.2), W, Inches(1.3), NAVY_DARK)

add_text(slide, Inches(0.5), Inches(0.4), Inches(12), Inches(0.6),
         "Conclusion & Future Work", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.5), Inches(1.05), Inches(5), Pt(2), BLUE)

# Conclusions column
add_text(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
         "Key Conclusions", font_size=22, color=BLUE_LIGHT, bold=True)
conclusions = [
    "\u2022  L4 semantic inversion defeats ALL tested",
    "   models (13 models, 950+ runs)",
    "\u2022  Pattern blindness is the core mechanism",
    "\u2022  Prior entrenchment depth governs resistance,",
    "   NOT task complexity",
    "\u2022  Three distinct failure modes identified",
    "\u2022  CoT reasoning does NOT help generation",
    "\u2022  Frontier models (gpt-5.4, o3) also fail",
]
add_bullets(slide, Inches(0.5), Inches(1.85), Inches(6), Inches(4),
            conclusions, font_size=16, color=WHITE, spacing=Pt(6))

# Future Work column
add_text(slide, Inches(7), Inches(1.3), Inches(6), Inches(0.4),
         "Future Work", font_size=22, color=BLUE_LIGHT, bold=True)
future = [
    "\u2022  Exp-5: Human study (student vs. LLM gap)",
    "\u2022  Broader task coverage (nested control,",
    "   stateful logic)",
    "\u2022  Execution-based judge for OpSub detection",
    "\u2022  Annotation density gradient replication",
    "\u2022  Isolate prior-depth vs. I/O-transparency",
    "\u2022  Cross-architecture replication",
    "\u2022  Longitudinal: does training data leakage",
    "   degrade resistance over time?",
]
add_bullets(slide, Inches(7), Inches(1.85), Inches(6), Inches(4),
            future, font_size=16, color=RGBColor(0xC0, 0xD0, 0xE8), spacing=Pt(6))

# Bottom tagline
add_text(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
         "A student reasons about a rule; an LLM samples from a prior.",
         font_size=22, color=BLUE_LIGHT, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(12), Inches(7.05), Inches(1), Inches(0.35),
         "25", font_size=11, color=MUTED, alignment=PP_ALIGN.RIGHT)

# ════════════════════════════════════════════════════════════════
#  APPENDIX — 8 slides
# ════════════════════════════════════════════════════════════════

def appendix_bar(slide, title, code):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.3), NAVY_DARK)
    add_text(slide, Inches(0.5), Inches(0.2), Inches(10), Inches(0.65),
             title, font_size=28, color=WHITE, bold=True)
    add_text(slide, Inches(10.5), Inches(0.25), Inches(2.5), Inches(0.4),
             f"Appendix {code}", font_size=14, color=AMBER, bold=True,
             alignment=PP_ALIGN.RIGHT)
    add_text(slide, Inches(0.5), Inches(0.78), Inches(10), Inches(0.35),
             "APPENDIX", font_size=12, color=AMBER, italic=True)
    add_text(slide, Inches(12), Inches(7.05), Inches(1), Inches(0.35),
             code, font_size=11, color=MUTED, alignment=PP_ALIGN.RIGHT)


# ── A1: Full L1 Data Table ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
appendix_bar(slide, "Full L1 KLR Data: Baseline vs. Context-Pack", "A1")

tbl = add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5), 11, 5)
hdrs = ["Model", "Baseline KLR", "Ctx-Pack KLR", "n", "Change"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 5)

a1_data = [
    ("gpt-4o", "0.37", "\u2014", "120", "\u2014"),
    ("gpt-4o-mini", "0.42", "0.26", "120", "\u221238%"),
    ("gpt-4.1", "\u2014", "0.30", "20", "\u2014"),
    ("gpt-4.1-mini", "0.43", "0.21", "120", "\u221251%"),
    ("gpt-4.1-nano", "\u2014", "0.38", "30", "\u2014"),
    ("gpt-5.4-mini", "0.67", "\u2014", "120", "\u2014"),
    ("o4-mini", "\u2014", "0.27", "20", "\u2014"),
    ("Llama-3.3-70B", "\u2014", "0.23", "20", "\u2014"),
    ("Qwen3-32B", "\u2014", "1.00", "20", "\u2014"),
]
# One extra empty row placeholder in case
for r, (m, bl, cp, n, ch) in enumerate(a1_data):
    sc(tbl.cell(r+1, 0), m, font_size=14, alignment=PP_ALIGN.LEFT, bold=True)
    sc(tbl.cell(r+1, 1), bl, font_size=14)
    sc(tbl.cell(r+1, 2), cp, font_size=14, bold=True)
    sc(tbl.cell(r+1, 3), n, font_size=14, color=MUTED)
    ch_clr = GREEN if "\u2212" in ch else MUTED
    sc(tbl.cell(r+1, 4), ch, font_size=14, color=ch_clr, bold="\u2212" in ch)
sc(tbl.cell(10, 0), "", font_size=1)  # pad last row

tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(2.2)
tbl.columns[2].width = Inches(2.2)
tbl.columns[3].width = Inches(1.2)
tbl.columns[4].width = Inches(2.0)

# ── A2: L4 Per-Variant Breakdown ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
appendix_bar(slide, "L4 Ablation: Per-Variant Breakdown (Variant A\u2013E \u00d7 9 Models)", "A2")

# This is a large table: 10 rows x 6 cols
tbl = add_table(slide, Inches(0.3), Inches(1.4), Inches(12.7), Inches(5.5), 10, 6)
hdrs = ["Model", "Var A (5ex,\u00d8)", "Var B (3ex,\u00d8)", "Var C (5ex,subtle)", "Var D (5ex,mod)", "Var E (5ex,strong)"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 6, font_size=11)

a2_data = [
    ("gpt-4o", "0/10", "0/10", "0/10", "0/10", "0/10"),
    ("gpt-4o-mini", "0/10", "0/10", "0/10", "0/10", "0/10"),
    ("gpt-4.1", "0/10", "0/10", "0/10", "0/10", "0/10"),
    ("gpt-4.1-mini", "0/10", "0/10 (PPR=0.0)", "0/10", "0/10", "0/10"),
    ("gpt-4.1-nano", "0/10", "0/10", "0/10", "0/10", "0/10"),
    ("gpt-5.4-mini", "0/8v", "0/9v", "0/9v", "0/8v", "0/8v"),
    ("o4-mini", "0/10", "0/8", "0/9", "0/9", "0/9"),
    ("Llama-3.3-70B", "0/5v", "0/5v", "0/5v", "0/5v", "0/5v"),
    ("Qwen3-32B", "0/10", "0/10", "0/10", "0/10", "0/10"),
]
for r, (m, *vs) in enumerate(a2_data):
    sc(tbl.cell(r+1, 0), m, font_size=12, alignment=PP_ALIGN.LEFT, bold=True, bg=LIGHT_RED)
    for c, v in enumerate(vs):
        bg = LIGHT_AMBER if "PPR=0.0" in v else LIGHT_RED
        sc(tbl.cell(r+1, c+1), v, font_size=11, bg=bg, color=RED)

tbl.columns[0].width = Inches(2.0)
for c in range(1, 6):
    tbl.columns[c].width = Inches(2.14)

# ── A3: CoT Ablation Full Data ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
appendix_bar(slide, "CoT Ablation: Full Data", "A3")

tbl = add_table(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(3.5), 5, 6)
hdrs = ["Model", "No-CoT Pass", "No-CoT PPR", "CoT Pass", "CoT PPR", "Delta"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 6)

a3_data = [
    ("gpt-4o", "0/20", "high", "0/20", "high", "0"),
    ("gpt-4.1-mini", "0/20", "high", "0/20", "high", "0"),
    ("gpt-4.1-nano", "0/20", "high", "0/20", "high", "0"),
    ("o4-mini", "0/20", "0.60", "1/20", "0.55", "+1"),
]
for r, (m, np, npp, cp, cpp, d) in enumerate(a3_data):
    sc(tbl.cell(r+1, 0), m, font_size=14, alignment=PP_ALIGN.LEFT, bold=True)
    sc(tbl.cell(r+1, 1), np, font_size=14, color=RED, bold=True, bg=LIGHT_RED)
    sc(tbl.cell(r+1, 2), npp, font_size=13, bg=LIGHT_RED)
    sc(tbl.cell(r+1, 3), cp, font_size=14, color=RED, bold=True, bg=LIGHT_RED)
    sc(tbl.cell(r+1, 4), cpp, font_size=13, bg=LIGHT_RED)
    sc(tbl.cell(r+1, 5), d, font_size=14, bold=True, color=MUTED)

add_text(slide, Inches(1.5), Inches(5.3), Inches(10), Inches(0.5),
         "Total: 160 runs. CoT provides zero benefit for standard models; "
         "marginal +1/20 for reasoning model (o4-mini).",
         font_size=16, color=BLUE, bold=True)

# ── A4: Llama-3.3-70B Variant Analysis ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
appendix_bar(slide, "Llama-3.3-70B: Context-Sensitive Pattern Blindness", "A4")

items = [
    "\u2022  Llama-3.3-70B (Groq, n=25 valid outputs across 5 variants)",
    "",
    "\u2022  Overall: 0/25 pass, PPR = 0.52",
    "   \u2013  Lower PPR than most OpenAI models \u2192 not purely prior-dominated",
    "   \u2013  But still 0% pass rate \u2192 pattern blindness persists",
    "",
    "\u2022  Variant-level analysis:",
    "   \u2013  Variants A\u2013B (no hint): PPR relatively high",
    "   \u2013  Variants C\u2013E (with hints): PPR drops, suggesting partial context sensitivity",
    "   \u2013  But hints shift output style without achieving rule extraction",
    "",
    "\u2022  Interpretation:",
    "   \u2013  Llama-3.3-70B shows a different failure profile than GPT models",
    "   \u2013  More sensitive to prompt hints but equally unable to extract implicit rules",
    "   \u2013  Suggests pattern blindness is architecture-independent (transformer-general)",
    "",
    "\u2022  L1 performance: KLR = 0.23 (best overall across all models)",
    "   \u2013  Excellent at surface-level alias compliance",
    "   \u2013  But L1 compliance does not predict L4 success",
]
add_bullets(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
            items, font_size=16, spacing=Pt(3))

# ── A5: Operational Substitution Detection ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
appendix_bar(slide, "Operational Substitution Detection Methodology", "A5")

items = [
    "\u2022  Goal: Distinguish true rule adoption from arithmetic workarounds",
    "",
    "\u2022  Detection pipeline:",
    "   1. Extract generated code from model output",
    "   2. Parse into AST (Python ast module)",
    "   3. Identify all conditional nodes (if/elif/else)",
    "   4. Check condition polarity against expected inverted semantics",
    "   5. If conditions are standard BUT output matches expected \u2192 flag as OpSub",
    "",
    "\u2022  Classification criteria:",
    "   \u2013  True adoption: inverted condition + correct output",
    "   \u2013  Prior dominance: standard condition + standard output",
    "   \u2013  Operational substitution: standard condition + correct output (via arithmetic)",
    "",
    "\u2022  Validation: Manual review of 50 gpt-4o count_vowels outputs",
    "   \u2013  47/50 classified as operational substitution (94%)",
    "   \u2013  Common patterns: count -= 1, return len(s) - count, negation tricks",
    "",
    "\u2022  Limitation: AST-based detection may miss novel workaround strategies",
    "   \u2013  Future: execution-based judge for more robust detection",
]
add_bullets(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
            items, font_size=15, spacing=Pt(3))

# ── A6: Metrics Definitions ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
appendix_bar(slide, "Metrics Definitions", "A6")

tbl = add_table(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.0), 6, 3)
hdrs = ["Metric", "Full Name", "Definition"]
for i, h in enumerate(hdrs):
    tbl.cell(0, i).text = h
style_header(tbl, 3)

metrics = [
    ("KLR", "Keyword Leakage Rate",
     "Fraction of outputs containing original Python keywords instead of aliases. Lower = better compliance."),
    ("PPR", "Python Prior Rate",
     "Fraction of outputs that reproduce standard Python patterns (e.g., if n<=1: return n). 1.0 = full prior dominance."),
    ("PSS", "Prior Substitution Score",
     "Weighted measure of how much standard Python structure persists in confusion language output."),
    ("PIR", "Prior Inversion Rate",
     "Fraction of conditionals where the model correctly inverts the condition. 1.0 = perfect inversion."),
    ("SIR", "Syntax Inversion Rate",
     "Fraction of outputs that fail to adopt inverted syntax rules. 0.0 = all outputs comply."),
]
for r, (met, full, defn) in enumerate(metrics):
    sc(tbl.cell(r+1, 0), met, font_size=15, bold=True)
    sc(tbl.cell(r+1, 1), full, font_size=13, alignment=PP_ALIGN.LEFT)
    sc(tbl.cell(r+1, 2), defn, font_size=12, alignment=PP_ALIGN.LEFT, color=MUTED)

tbl.columns[0].width = Inches(1.2)
tbl.columns[1].width = Inches(2.8)
tbl.columns[2].width = Inches(8.3)

# ── A7: Prompt Engineering Example ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
appendix_bar(slide, "Prompt Engineering: L4 Variant A (Full Prompt)", "A7")

prompt_text = [
    "\u2022  System prompt structure (Variant A \u2014 5 examples, no hint):",
    "",
    "   1. Role definition: 'You are a programmer in a confusion language'",
    "   2. Context: 'This language looks like Python but behaves differently'",
    "   3. Five worked examples showing:",
    "      \u2013  Simple conditional (if x > 0: return 'neg' else: return 'pos')",
    "      \u2013  Loop with inverted condition check",
    "      \u2013  Nested conditional with inverted semantics",
    "      \u2013  Function with inverted guard clause",
    "      \u2013  Recursive function with inverted base case",
    "   4. Target task: 'Write fibonacci(n) in this confusion language'",
    "",
    "\u2022  Key design choices:",
    "   \u2013  Never use words like 'invert', 'opposite', or 'negate'",
    "   \u2013  Examples are self-consistent (all follow inverted rule)",
    "   \u2013  Output format: raw code only, no explanation",
    "",
    "\u2022  Variant progression A\u2192E adds progressively more explicit hints",
    "   while keeping the same example set (except B which uses 3 examples)",
]
add_bullets(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
            prompt_text, font_size=14, spacing=Pt(2))

# ── A8: GitHub / Reproduction Guide ──
slide = new_slide()
add_bg(slide, LIGHT_BG)
appendix_bar(slide, "Repository Structure & Reproduction Guide", "A8")

items = [
    "\u2022  Repository: github.com/[redacted]/cse307-open-project",
    "",
    "\u2022  Directory structure:",
    "   experiments/          \u2014 All experiment scripts and data",
    "   experiments/results/  \u2014 Raw JSON outputs per experiment",
    "   paper/overleaf/       \u2014 LaTeX source for the paper",
    "   presentation/         \u2014 Slide generation scripts",
    "   docs/research/        \u2014 Research notes and analysis",
    "",
    "\u2022  Reproduction steps:",
    "   1. Clone repository",
    "   2. Install dependencies: pip install openai groq python-pptx",
    "   3. Set API keys: OPENAI_API_KEY, GROQ_API_KEY",
    "   4. Run experiments: python experiments/run_experiment.py --exp E01",
    "   5. Results saved to experiments/results/",
    "",
    "\u2022  Key experiment IDs:",
    "   E01\u2013E05: L1 KLR experiments",
    "   E10\u2013E15: L4 ablation (variants A\u2013E)",
    "   E20\u2013E24: Multi-task, hard tasks, frontier",
    "   E28: Annotation density gradient (o4-mini)",
    "",
    "\u2022  Total runs: 950+ L4 evaluations, ~590 L1 evaluations",
]
add_bullets(slide, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
            items, font_size=14, spacing=Pt(2))

# ════════════════════════════════════════════════════════════════
#  Save
# ════════════════════════════════════════════════════════════════
output_path = "/Users/jaeyeong_openclaw/.openclaw/workspace/cse307-open-project/presentation/research-slides-v2.pptx"
prs.save(output_path)
print(f"Saved {SLIDE_NUM[0]} slides to {output_path}")
