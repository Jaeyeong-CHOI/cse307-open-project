#!/usr/bin/env python3
"""Build research presentation slides using python-pptx."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
NAVY = RGBColor(0x1A, 0x27, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF6, 0xFA)
BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_TEXT = RGBColor(0x6B, 0x7B, 0x8D)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)
LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK_TEXT,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_table(slide, left, top, width, height, rows, cols):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def style_header_row(table, col_count, bg=NAVY, fg=WHITE):
    for i in range(col_count):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = fg
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER


def style_data_cell(cell, text, font_size=13, color=DARK_TEXT, bold=False,
                    alignment=PP_ALIGN.CENTER, bg=None):
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = alignment
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg


def slide_title_bar(slide, title_text, subtitle_text=None):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.4), NAVY)
    add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
             title_text, font_size=32, color=WHITE, bold=True)
    if subtitle_text:
        add_text(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
                 subtitle_text, font_size=16, color=RGBColor(0xA0, 0xB0, 0xC8))


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, NAVY)

add_text(slide, Inches(1.5), Inches(1.2), Inches(10.3), Inches(1.5),
         "Inverting Python:", font_size=44, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.2),
         "Designing NTP-Prior-Resistant Programming Languages\ntoward LLM-Resistant Coding Assessment",
         font_size=24, color=RGBColor(0xA0, 0xC4, 0xF0), alignment=PP_ALIGN.CENTER)

# Divider line
add_rect(slide, Inches(4.5), Inches(3.8), Inches(4.3), Pt(2), BLUE)

add_text(slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
         "Jaeyeong CHOI", font_size=22, color=WHITE, bold=True,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.5),
         "DGIST EECS  |  jaeyeong2022@dgist.ac.kr", font_size=16,
         color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.5),
         "13 Models  |  950+ L4 Evaluations  |  4-Level Taxonomy",
         font_size=16, color=RGBColor(0x7B, 0x9B, 0xBB), alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Research Motivation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_title_bar(slide, "Research Motivation", "Why This Research?")

items = [
    "LLMs now solve standard programming tasks well enough to undermine assessment value",
    "Students can use ChatGPT/Copilot to pass coding assignments without understanding",
    "Traditional countermeasures (plagiarism detection, proctoring) are insufficient for LLM-generated code",
    "Core Question: Can we design a language that humans can learn but LLMs cannot solve?",
    "Key Insight: LLM code generation is governed by NTP (Next-Token Prediction) priors from training data",
    "If we invert the semantics that LLMs have deeply memorized, we create a natural human-LLM asymmetry",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5),
                [f"\u2022  {item}" for item in items], font_size=20, spacing=Pt(12))

# ============================================================
# SLIDE 3: Core Asymmetry
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_title_bar(slide, "Core Asymmetry", "Human vs. LLM: Fundamentally Different Mechanisms")

# Human box
human_box = add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), WHITE, BLUE)
add_text(slide, Inches(1.2), Inches(1.95), Inches(4.8), Inches(0.5),
         "Human Student", font_size=26, color=BLUE, bold=True)
items_h = [
    "\u2022  Reads explicit rule from instructor",
    "\u2022  Verifies rule against worked examples",
    "\u2022  Applies rule step-by-step",
    "\u2022  Reasoning-based: rule \u2192 application",
    "\u2022  Can override prior knowledge with new rules",
]
add_bullet_text(slide, Inches(1.2), Inches(2.7), Inches(4.8), Inches(3),
                items_h, font_size=17, spacing=Pt(8))
add_text(slide, Inches(1.2), Inches(5.2), Inches(4.8), Inches(0.5),
         "Result: Can learn inverted semantics", font_size=15, color=GREEN, bold=True)

# LLM box
llm_box = add_rect(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), WHITE, RED)
add_text(slide, Inches(7.4), Inches(1.95), Inches(4.8), Inches(0.5),
         "LLM (GPT-4o, o4-mini, etc.)", font_size=26, color=RED, bold=True)
items_l = [
    "\u2022  Samples from NTP prior distribution",
    "\u2022  Prior built from billions of Python tokens",
    "\u2022  if n <= 1: return n is deeply entrenched",
    "\u2022  Generation-based: prior \u2192 sampling",
    "\u2022  Cannot override deep priors via in-context cues",
]
add_bullet_text(slide, Inches(7.4), Inches(2.7), Inches(4.8), Inches(3),
                items_l, font_size=17, spacing=Pt(8))
add_text(slide, Inches(7.4), Inches(5.2), Inches(4.8), Inches(0.5),
         "Result: Fails on inverted Fibonacci (0%)", font_size=15, color=RED, bold=True)

# ============================================================
# SLIDE 4: 4-Level Taxonomy
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_title_bar(slide, "4-Level Confusion Language Taxonomy", "Perturbation Depth x Rule Delivery")

tbl = add_table(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(4.2), 5, 5)
headers = ["Level", "Perturbation", "Rule Delivery", "Observed Result", "Failure Mode"]
for i, h in enumerate(headers):
    tbl.cell(0, i).text = h
style_header_row(tbl, 5)

data = [
    ["L1", "Token Alias", "Explicit", "Partial resistance (KLR 0.21\u20130.69)", "Prior Dominance"],
    ["L2", "Syntax Inversion", "Explicit", "Fully adopted by all models", "None (adopted)"],
    ["L3", "Semantic Inversion", "Explicit", "Capacity-dependent", "Model threshold"],
    ["L4", "Semantic Inversion", "Implicit (examples only)", "ALL models fail (0% pass)", "Pattern Blindness"],
]
colors_bg = [
    RGBColor(0xFF, 0xFB, 0xEB),  # yellow-ish
    RGBColor(0xFE, 0xE2, 0xE2),  # red-ish (not resistant)
    RGBColor(0xFF, 0xFB, 0xEB),  # yellow-ish
    LIGHT_GREEN,
]
for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        style_data_cell(tbl.cell(r + 1, c), val, font_size=14, bg=colors_bg[r])

# Set column widths
col_widths = [Inches(0.9), Inches(2.2), Inches(2.2), Inches(3.8), Inches(2.6)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

add_text(slide, Inches(0.8), Inches(6.1), Inches(11), Inches(0.5),
         "L4 is the only level where ALL tested models fail regardless of capacity (950+ evaluations, 13 models)",
         font_size=16, color=BLUE, bold=True)

# ============================================================
# SLIDE 5: Experimental Setup
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_title_bar(slide, "Experimental Setup", "Models, Tasks, and Evaluation Protocol")

# Models box
add_rect(slide, Inches(0.6), Inches(1.7), Inches(5.8), Inches(2.6), WHITE, BLUE)
add_text(slide, Inches(0.9), Inches(1.8), Inches(5), Inches(0.4),
         "13 Models Tested", font_size=20, color=BLUE, bold=True)
model_items = [
    "\u2022  OpenAI Core (7): gpt-4o, gpt-4o-mini, gpt-4.1, gpt-4.1-mini, gpt-4.1-nano, gpt-5.4-mini, o4-mini",
    "\u2022  Frontier (2): gpt-5.4, o3",
    "\u2022  Open-weight (2): Llama-3.3-70B, Qwen3-32B",
]
add_bullet_text(slide, Inches(0.9), Inches(2.3), Inches(5.3), Inches(2),
                model_items, font_size=14, spacing=Pt(6))

# Tasks box
add_rect(slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(2.6), WHITE, BLUE)
add_text(slide, Inches(7.1), Inches(1.8), Inches(5), Inches(0.4),
         "Evaluation Tasks", font_size=20, color=BLUE, bold=True)
task_items = [
    "\u2022  T1: Fibonacci (deep prior entrenchment)",
    "\u2022  T2: is_sorted (shallow prior)",
    "\u2022  T3: count_vowels (operational substitution target)",
    "\u2022  H1-H3: merge_sort, binary_search, BFS (hard tasks)",
]
add_bullet_text(slide, Inches(7.1), Inches(2.3), Inches(5.3), Inches(2),
                task_items, font_size=14, spacing=Pt(6))

# Protocol box
add_rect(slide, Inches(0.6), Inches(4.6), Inches(12), Inches(2.4), WHITE, BLUE)
add_text(slide, Inches(0.9), Inches(4.7), Inches(5), Inches(0.4),
         "Evaluation Protocol", font_size=20, color=BLUE, bold=True)
proto_items = [
    "\u2022  950+ total L4 evaluations across all conditions",
    "\u2022  L4 ablation: 7 models x 5 variants x n=50 = 350 runs (strict example-only)",
    "\u2022  Multi-task: 345 runs (3 tasks x 7 models, partially-annotated)",
    "\u2022  CoT ablation: 160 runs (4 models x 2 conditions x n=20)",
    "\u2022  Metrics: Pass Rate, PPR (Python Prior Rate), KLR (Keyword Leakage Rate)",
]
add_bullet_text(slide, Inches(0.9), Inches(5.2), Inches(11), Inches(1.8),
                proto_items, font_size=14, spacing=Pt(4))

# ============================================================
# SLIDE 6: Key Results - L1 KLR
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_title_bar(slide, "Key Results: L1 Keyword Leakage Rate", "Token Alias Substitution Creates Only Partial Resistance")

tbl = add_table(slide, Inches(0.8), Inches(1.8), Inches(7), Inches(4.5), 10, 3)
headers = ["Model", "Delivery", "KLR"]
for i, h in enumerate(headers):
    tbl.cell(0, i).text = h
style_header_row(tbl, 3)

klr_data = [
    ("gpt-4.1-mini", "ctx-pack", "0.21", LIGHT_GREEN),
    ("Llama-3.3-70B", "ctx-pack", "0.23", LIGHT_GREEN),
    ("gpt-4o-mini", "ctx-pack", "0.26", LIGHT_GREEN),
    ("o4-mini", "ctx-pack", "0.27", LIGHT_GREEN),
    ("gpt-4.1", "ctx-pack", "0.30", None),
    ("gpt-4o", "baseline", "0.37", None),
    ("gpt-5.4-mini", "baseline", "0.67", LIGHT_RED),
    ("gpt-4.1-nano", "baseline", "0.69", LIGHT_RED),
    ("Qwen3-32B", "ctx-pack", "1.00", LIGHT_RED),
]
for r, (model, delivery, klr, bg) in enumerate(klr_data):
    style_data_cell(tbl.cell(r + 1, 0), model, font_size=13, alignment=PP_ALIGN.LEFT, bg=bg)
    style_data_cell(tbl.cell(r + 1, 1), delivery, font_size=13, bg=bg)
    style_data_cell(tbl.cell(r + 1, 2), klr, font_size=13, bold=True, bg=bg)

tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(1.8)
tbl.columns[2].width = Inches(1.2)

# Key findings
findings = [
    "\u2022  Context-pack reduces KLR: gpt-4o-mini 0.42\u21920.26 (-38%)",
    "\u2022  gpt-5.4-mini (KLR=0.67) worse than older gpt-4o (0.37)",
    "\u2022  Qwen3-32B: worst result (KLR=1.00) across all models",
    "\u2022  Cross-family variance (0.23 vs 1.00) > intra-family",
    "\u2022  Architecture & fine-tuning govern alias compliance",
]
add_rect(slide, Inches(8.3), Inches(1.8), Inches(4.5), Inches(4.5), WHITE, BLUE)
add_text(slide, Inches(8.5), Inches(1.9), Inches(4), Inches(0.4),
         "Key Findings", font_size=18, color=BLUE, bold=True)
add_bullet_text(slide, Inches(8.5), Inches(2.5), Inches(4.1), Inches(3.5),
                findings, font_size=14, spacing=Pt(8))

# ============================================================
# SLIDE 7: Key Results - L4 Pattern Blindness
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_title_bar(slide, "Key Results: L4 Pattern Blindness",
                "ALL Models Fail on Prior-Entrenched Tasks (0% Pass Rate)")

tbl = add_table(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(2.8), 4, 8)
headers2 = ["", "gpt-4o", "gpt-4o-mini", "gpt-4.1", "gpt-4.1-mini", "gpt-4.1-nano", "gpt-5.4-mini", "o4-mini"]
for i, h in enumerate(headers2):
    tbl.cell(0, i).text = h
style_header_row(tbl, 8)

style_data_cell(tbl.cell(1, 0), "Pass Rate", bold=True, alignment=PP_ALIGN.LEFT)
style_data_cell(tbl.cell(2, 0), "PPR", bold=True, alignment=PP_ALIGN.LEFT)
style_data_cell(tbl.cell(3, 0), "n", bold=True, alignment=PP_ALIGN.LEFT)

pass_rates = ["0/50", "0/50", "0/50", "0/50", "0/50", "0/50", "0/50"]
pprs = ["0.70", "0.82", "1.00", "0.38", "1.00", "1.00", "0.92"]
ns = ["50", "50", "50", "50", "50", "50", "50"]
for i, (pr, ppr, n) in enumerate(zip(pass_rates, pprs, ns)):
    style_data_cell(tbl.cell(1, i + 1), pr, color=RED, bold=True, bg=LIGHT_RED)
    style_data_cell(tbl.cell(2, i + 1), ppr, font_size=13)
    style_data_cell(tbl.cell(3, i + 1), n, font_size=12, color=GRAY_TEXT)

tbl.columns[0].width = Inches(1.4)
for i in range(1, 8):
    tbl.columns[i].width = Inches(1.56)

# E24 box
add_rect(slide, Inches(0.5), Inches(4.7), Inches(6), Inches(2.3), WHITE, RED)
add_text(slide, Inches(0.8), Inches(4.85), Inches(5.5), Inches(0.4),
         "E24: Frontier Models Also Fail", font_size=20, color=RED, bold=True)
frontier_items = [
    "\u2022  gpt-5.4:  0/10  (PPR = 1.00)",
    "\u2022  o3:         0/10  (PPR = 1.00)",
    "",
    "Pattern blindness persists even in the",
    "latest frontier models (as of 2026-03-25)",
]
add_bullet_text(slide, Inches(0.8), Inches(5.4), Inches(5.5), Inches(1.5),
                frontier_items, font_size=16, spacing=Pt(4))

# CoT box
add_rect(slide, Inches(6.8), Inches(4.7), Inches(6), Inches(2.3), WHITE, BLUE)
add_text(slide, Inches(7.1), Inches(4.85), Inches(5.5), Inches(0.4),
         "CoT Ablation: Reasoning \u2260 Generation", font_size=20, color=BLUE, bold=True)
cot_items = [
    "\u2022  160 CoT runs across 4 models",
    "\u2022  Models articulate inversion rule (100% mentions)",
    "\u2022  But still generate standard Python code",
    "\u2022  Reasoning-generation dissociation confirmed",
]
add_bullet_text(slide, Inches(7.1), Inches(5.4), Inches(5.5), Inches(1.5),
                cot_items, font_size=16, spacing=Pt(4))

# ============================================================
# SLIDE 8: Hard Tasks
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_title_bar(slide, "New Finding: Hard Tasks vs. Prior Entrenchment",
                "Task Complexity \u2260 LLM Resistance")

tbl = add_table(slide, Inches(0.8), Inches(1.8), Inches(8), Inches(3), 5, 5)
headers3 = ["Model", "H1: merge_sort", "H2: bin_search", "H3: BFS", "Overall"]
for i, h in enumerate(headers3):
    tbl.cell(0, i).text = h
style_header_row(tbl, 5)

hard_data = [
    ("gpt-4o", "4/10", "10/10", "10/10", "80%"),
    ("gpt-5.4", "10/10", "10/10", "10/10", "100%"),
    ("o3", "10/10", "10/10", "10/10", "100%"),
    ("o4-mini", "10/10", "10/10", "10/10", "100%"),
]
for r, (model, h1, h2, h3, overall) in enumerate(hard_data):
    style_data_cell(tbl.cell(r + 1, 0), model, alignment=PP_ALIGN.LEFT, bold=True)
    h1_bg = RGBColor(0xFF, 0xFB, 0xEB) if model == "gpt-4o" else LIGHT_GREEN
    style_data_cell(tbl.cell(r + 1, 1), h1, bg=h1_bg, color=DARK_TEXT if model == "gpt-4o" else GREEN)
    style_data_cell(tbl.cell(r + 1, 2), h2, bg=LIGHT_GREEN, color=GREEN)
    style_data_cell(tbl.cell(r + 1, 3), h3, bg=LIGHT_GREEN, color=GREEN)
    ov_color = DARK_TEXT if model == "gpt-4o" else GREEN
    style_data_cell(tbl.cell(r + 1, 4), overall, bold=True, color=ov_color)

# Contrast box
add_rect(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(2), WHITE, BLUE)
add_text(slide, Inches(1.1), Inches(5.2), Inches(5), Inches(0.4),
         "Contrast with L4 Fibonacci", font_size=18, color=BLUE, bold=True)
contrast = [
    "\u2022  Same models on Fibonacci: ALL fail (0%)",
    "\u2022  Same models on hard tasks: near 100%",
    "\u2022  Prior entrenchment depth, not task difficulty,",
    "   governs LLM resistance",
]
add_bullet_text(slide, Inches(1.1), Inches(5.7), Inches(5), Inches(1.3),
                contrast, font_size=15, spacing=Pt(4))

# Implication box
add_rect(slide, Inches(6.8), Inches(5.1), Inches(5.8), Inches(2), WHITE, RED)
add_text(slide, Inches(7.1), Inches(5.2), Inches(5.3), Inches(0.4),
         "Key Implication", font_size=18, color=RED, bold=True)
imp = [
    "\u2022  Harder \u2260 more LLM-resistant",
    "\u2022  Deep, uniform priors (Fibonacci) = resistant",
    "\u2022  Varied implementations (BFS, sort) = extractable",
    "\u2022  Task selection is critical for assessment design",
]
add_bullet_text(slide, Inches(7.1), Inches(5.7), Inches(5.3), Inches(1.3),
                imp, font_size=15, spacing=Pt(4))

# ============================================================
# SLIDE 9: 3 Failure Modes
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_title_bar(slide, "Three Failure Modes", "Distinct Mechanisms Behind LLM Failure on Confusion Languages")

# Mode 1
add_rect(slide, Inches(0.5), Inches(1.8), Inches(3.9), Inches(4.8), WHITE, BLUE)
add_text(slide, Inches(0.8), Inches(1.95), Inches(3.4), Inches(0.4),
         "1. Prior Dominance", font_size=22, color=BLUE, bold=True)
pd_items = [
    "Stated or implicit rule is overridden by pretrained token distribution",
    "",
    "Where: L1, L3-weak, L4-T1",
    "",
    "Example: Model generates if n <= 1: return n despite being told if-blocks run when FALSE",
    "",
    "83.1% of all L4 failures (Type-I)",
]
add_bullet_text(slide, Inches(0.8), Inches(2.6), Inches(3.4), Inches(3.8),
                pd_items, font_size=13, spacing=Pt(2))

# Mode 2
add_rect(slide, Inches(4.7), Inches(1.8), Inches(3.9), Inches(4.8), WHITE, RED)
add_text(slide, Inches(5.0), Inches(1.95), Inches(3.4), Inches(0.4),
         "2. Pattern Blindness", font_size=22, color=RED, bold=True)
pb_items = [
    "Model cannot extract semantic rule from examples alone",
    "",
    "Where: L4-T1 (Fibonacci)",
    "",
    "Key evidence: Same model follows rule at L3 (explicit) but fails at L4 (implicit)",
    "",
    "Universal across all 13 tested models on Fibonacci",
]
add_bullet_text(slide, Inches(5.0), Inches(2.6), Inches(3.4), Inches(3.8),
                pb_items, font_size=13, spacing=Pt(2))

# Mode 3
add_rect(slide, Inches(8.9), Inches(1.8), Inches(3.9), Inches(4.8), WHITE, RGBColor(0xEA, 0x58, 0x0C))
add_text(slide, Inches(9.2), Inches(1.95), Inches(3.4), Inches(0.4),
         "3. Operational Substitution", font_size=22, color=RGBColor(0xEA, 0x58, 0x0C), bold=True)
os_items = [
    "Model compensates via arithmetic transformation instead of adopting rule",
    "",
    "Where: L4-T3 (count_vowels)",
    "",
    "Example: count -= 1; return -count achieves correct output without rule",
    "",
    "94% of gpt-4o T3 runs (47/50, CI: [83%, 99%])",
]
add_bullet_text(slide, Inches(9.2), Inches(2.6), Inches(3.4), Inches(3.8),
                os_items, font_size=13, spacing=Pt(2))

# ============================================================
# SLIDE 10: Educational Implications
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
slide_title_bar(slide, "Educational Implications",
                "Practical Guidance for LLM-Resistant Programming Assessment")

items_edu = [
    "\u2022  L4 confusion languages restore discriminative power in coding assessments",
    "\u2022  Students who understand rules (from lecture) retain advantage over LLM-assisted submissions",
    "\u2022  Current models universally fail L4 Fibonacci while humans can apply rules from examples",
    "",
    "\u2022  Design Recipe for Instructors:",
    "     1. Retain Python surface syntax (students already know it)",
    "     2. Invert conditional semantics (if-blocks execute when condition is FALSE)",
    "     3. Provide only worked examples, NOT explicit rules",
    "     4. Choose tasks with deep Python priors (e.g., Fibonacci recursion)",
    "",
    "\u2022  Caveats:",
    "     \u2013  Tasks with shallow priors (is_sorted) may still be solvable by capable models",
    "     \u2013  Partial annotation can help reasoning models (o4-mini: 8/15 with hints)",
    "     \u2013  Human baseline study (Exp-5) is needed to quantify the exact performance gap",
]
add_bullet_text(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5.2),
                items_edu, font_size=17, spacing=Pt(4))

# ============================================================
# SLIDE 11: Conclusion & Future Work
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)

add_text(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.7),
         "Conclusion & Future Work", font_size=36, color=WHITE, bold=True)
add_rect(slide, Inches(0.8), Inches(1.2), Inches(5), Pt(2), BLUE)

# Conclusions
add_text(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
         "Key Conclusions", font_size=22, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)
conclusions = [
    "\u2022  L4 semantic inversion defeats ALL tested models (13 models, 950+ runs)",
    "\u2022  Pattern blindness is the core mechanism: models cannot extract implicit rules",
    "\u2022  Prior entrenchment depth governs resistance, NOT task complexity",
    "\u2022  Three distinct failure modes identified and quantified",
    "\u2022  CoT reasoning helps articulate but NOT generate correct inverted code",
    "\u2022  Frontier models (gpt-5.4, o3) also fail on prior-entrenched tasks",
]
add_bullet_text(slide, Inches(0.8), Inches(2.1), Inches(5.8), Inches(4),
                conclusions, font_size=15, color=WHITE, spacing=Pt(8))

# Future Work
add_text(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
         "Future Work", font_size=22, color=RGBColor(0x60, 0xA5, 0xFA), bold=True)
future = [
    "\u2022  Exp-5: Human study (student vs. LLM performance gap)",
    "\u2022  Broader task coverage (nested control-flow, stateful logic)",
    "\u2022  Execution-based judge for operational substitution",
    "\u2022  Annotation density gradient for other reasoning models",
    "\u2022  Isolate prior-depth vs. I/O-transparency confound",
    "\u2022  Cross-architecture replication (more open-weight models)",
]
add_bullet_text(slide, Inches(7), Inches(2.1), Inches(5.8), Inches(4),
                future, font_size=15, color=RGBColor(0xD0, 0xD8, 0xE8), spacing=Pt(8))

# Bottom tagline
add_rect(slide, Inches(0), Inches(6.6), W, Inches(0.9), RGBColor(0x10, 0x1B, 0x30))
add_text(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.5),
         "A student reasons about a rule; an LLM samples from a prior.",
         font_size=20, color=RGBColor(0x60, 0xA5, 0xFA), bold=True, alignment=PP_ALIGN.CENTER)

# Save
output_path = "/Users/jaeyeong_openclaw/.openclaw/workspace/cse307-open-project/presentation/research-slides.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
